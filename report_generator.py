import re
import os
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.action import PP_ACTION
from dataclasses import dataclass, field
from typing import List, Dict


@dataclass
class SlideData:
    slide_num: int
    content: Dict[str, str] = field(default_factory=dict)
    images: Dict[int, str] = field(default_factory=dict)

class PowerPointModifier:
    def __init__(self, template_path: str, output_path: str, csv_path: str):
        self.prs = Presentation(template_path)
        self.output_path = output_path
        self.slide_data = self.load_slide_data(csv_path)
        self.text_placeholder_pattern = re.compile(r'\{\{(.*?)\}\}', re.DOTALL)
        self.image_extensions = ('.png', '.jpg', '.jpeg', '.gif')

    def load_slide_data(self, csv_path: str) -> List[SlideData]:
        slide_data_dict = {}
        df = pd.read_csv(csv_path)
        for _, row in df.iterrows():
            slide_num = int(row['slide_num'])
            placeholder = row['placeholder']
            value = str(row['value']) if not pd.isnull(row['value']) else ''
            if slide_num not in slide_data_dict:
                slide_data_dict[slide_num] = SlideData(slide_num=slide_num)
            if os.path.isfile(value) and placeholder.isdigit():  
                placeholder_idx = int(placeholder)
                slide_data_dict[slide_num].images[placeholder_idx] = value
            slide_data_dict[slide_num].content[placeholder] = value
        return list(slide_data_dict.values())

    def print_placeholders(self):
        for slide_index, slide in enumerate(self.prs.slides):
            print(f"Slide {slide_index + 1}:")
            for shape in slide.shapes:
                if shape.is_placeholder:
                    placeholder = shape.placeholder_format
                    print(f"  Placeholder {placeholder.idx} - {placeholder.type}")

    def replace_placeholders(self):
        for slide_index, slide in enumerate(self.prs.slides):
            for data in self.slide_data:
                if slide_index + 1 == data.slide_num:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            self._replace_text_placeholders(slide, shape, data.content)
                        if shape.has_table:
                            self._replace_table_placeholder(shape, data.content)
                        if shape.is_placeholder and shape.placeholder_format.idx in data.images:
                            self._replace_image_placeholder(slide, shape, data.images[shape.placeholder_format.idx])

    def _replace_text_placeholders(self, slide, shape, content):
        text_frame = shape.text_frame
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                new_text = self.text_placeholder_pattern.sub(
                    lambda match: self._replace_text_or_image(slide, shape, match, content),
                    run.text
                )
                run.text = new_text
                self._identify_and_replace_hyperlink(slide, shape, content)

    def _replace_text_or_image(self, slide, shape, match, content):
        placeholder = match.group(1).strip()
        if any(placeholder.lower().endswith(ext) for ext in self.image_extensions):
            image_path = content.get(placeholder, match.group(0))
            if os.path.isfile(image_path):
                self._replace_image_placeholder(slide, shape, image_path)
                return ""
            return match.group(0)
        else:
            return str(content.get(placeholder, match.group(0)))

    def _replace_image_placeholder(self, slide, shape, image_path):
        left = shape.left
        top = shape.top
        width = shape.width
        height = shape.height
        slide.shapes.add_picture(image_path, left, top, width, height)
        shape.element.getparent().remove(shape.element)

    def _replace_table_placeholder(self, shape, content):
        table = shape.table
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        new_text = self.text_placeholder_pattern.sub(
                            lambda match: self._replace_text(match, content),
                            run.text
                        )
                        run.text = new_text

    def _replace_text(self, match, content):
        placeholder = match.group(1).strip()
        return str(content.get(placeholder, match.group(0)))

    def _identify_and_replace_hyperlink(self, slide, shape, content):
        text = shape.text
        for content_key, new_link in content.items():
            if content_key in text:  # Check if 'Content' is in the shape's text
                # Check if the shape has a click action with a hyperlink
                if shape.click_action and shape.click_action.action == PP_ACTION.HYPERLINK:
                    shape.click_action.hyperlink.address = new_link
                    #print(f"Slide {slide.slide_number}: Replaced hyperlink for '{content_key}' -> {new_link}")
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.hyperlink.address:
                            run.hyperlink.address = new_link
                            #print(f"Slide {slide.slide_number}: Replaced hyperlink for '{content_key}' -> {new_link}")

    def save_presentation(self):
        self.prs.save(self.output_path)

csv_path = 'input.csv'
ppt_modifier = PowerPointModifier(
    template_path='TEMPLATE1.pptx',
    output_path='template_output.pptx',
    csv_path=csv_path
)
ppt_modifier.print_placeholders()
ppt_modifier.replace_placeholders()
ppt_modifier.save_presentation()
